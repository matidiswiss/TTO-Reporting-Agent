#!/usr/bin/env python3
"""Generate preview/data.json and FOLO-style PowerPoint decks from TikTok One Excel exports."""

from __future__ import annotations

import argparse
import json
import re
import shutil
import sys
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import Any

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
DEFAULT_TEMPLATE = ROOT / "references" / "FOLO - Report TTO Wound (Mar-Apr).pptx"
DEFAULT_JSON = ROOT / "preview" / "data.json"
DEFAULT_OUTPUT_DIR = ROOT / "output"

GREEN = RGBColor(0x15, 0x81, 0x58)
GRAY = RGBColor(0x59, 0x59, 0x59)

REQUIRED_SHEETS = ("Campaign Report", "Campaign History Report", "Video Report")


@dataclass
class ReportContext:
    campaign: dict[str, Any]
    creators: list[dict[str, Any]]
    top_videos: list[dict[str, Any]]
    history: list[dict[str, Any]]
    meta: dict[str, Any]
    qa_notes: list[str]
    videos_df: pd.DataFrame
    period: str
    brand: str
    source: str
    product_label: str


def fmt_num(n: float | int) -> str:
    n = float(n)
    if n >= 1_000_000:
        s = f"{n / 1_000_000:.1f}M"
        return s.replace(".0M", "M")
    if n >= 1_000:
        s = f"{n / 1_000:.1f}K"
        return s.replace(".0K", "K")
    return f"{int(n):,}"


def pct(part: float, total: float) -> int:
    return round(part / total * 100) if total else 0


def json_default(obj: Any) -> Any:
    import numpy as np

    if isinstance(obj, (str, bool)):
        return obj
    if isinstance(obj, (np.integer, int)):
        return int(obj)
    if isinstance(obj, (np.floating, float)):
        return None if pd.isna(obj) else float(obj)
    if isinstance(obj, np.ndarray):
        return obj.tolist()
    if pd.isna(obj):
        return None
    if isinstance(obj, (datetime, pd.Timestamp)):
        return obj.isoformat()
    raise TypeError(f"Object of type {type(obj)} is not JSON serializable")


def row_to_dict(row: pd.Series) -> dict[str, Any]:
    return {k: json_default(v) for k, v in row.items()}


def parse_distribution(raw: Any, key: str) -> list[dict[str, Any]]:
    if raw is None or (isinstance(raw, float) and pd.isna(raw)):
        return []
    try:
        data = json.loads(str(raw).strip())
    except json.JSONDecodeError:
        return []
    return sorted(data, key=lambda item: item.get("percentage", 0), reverse=True)


def infer_period(history_df: pd.DataFrame) -> str:
    if history_df.empty:
        return "—"
    dates = pd.to_datetime(history_df["Date"].astype(str), format="%Y%m%d", errors="coerce")
    dates = dates.dropna()
    if dates.empty:
        return "—"
    start, end = dates.min(), dates.max()
    if start.year == end.year and start.month == end.month:
        return start.strftime("%b %Y")
    if start.year == end.year:
        return f"{start.strftime('%b')}–{end.strftime('%b %Y')}"
    return f"{start.strftime('%b %Y')}–{end.strftime('%b %Y')}"

def infer_product_label(campaign_name: str, input_path: Path) -> str:
    stem = input_path.stem
    if " - " in stem:
        return stem.split(" - ", 1)[1].strip()
    return str(campaign_name)


def clean_product_label(label: str) -> str:
    return re.sub(r"\s*\([^)]*OLV[^)]*\)\s*", "", label, flags=re.I).strip()


def wound_title_subtitle(brand: str) -> str:
    match = re.search(r"\(([^)]+)\)\s*$", brand)
    if match:
        return match.group(1)
    return brand


def default_brand() -> str:
    return "Betadine Wound (Solution, Ointment, Bening)"


def load_excel(
    input_path: Path,
    *,
    period: str | None,
    brand: str | None,
    source: str | None,
) -> ReportContext:
    if not input_path.exists():
        raise FileNotFoundError(f"Excel file not found: {input_path}")

    xl = pd.ExcelFile(input_path)
    missing = [sheet for sheet in REQUIRED_SHEETS if sheet not in xl.sheet_names]
    if missing:
        raise ValueError(f"Missing required sheet(s): {', '.join(missing)}")

    campaign_df = pd.read_excel(input_path, sheet_name="Campaign Report")
    history_df = pd.read_excel(input_path, sheet_name="Campaign History Report")
    videos_df = pd.read_excel(input_path, sheet_name="Video Report")

    campaign = row_to_dict(campaign_df.iloc[0])
    history_df = history_df.copy()
    history_df["Date"] = pd.to_datetime(
        history_df["Date"].astype(str), format="%Y%m%d", errors="coerce"
    )

    creators = (
        videos_df.groupby("Creator name", as_index=False)
        .agg(
            videos=("Video ID", "count"),
            total_views=("Total views", "sum"),
            paid_views=("Paid views", "sum"),
            organic_views=("Organic views", "sum"),
            likes=("Likes", "sum"),
            shares=("Shares", "sum"),
            engagement_rate=("Engagement rate", "mean"),
        )
        .sort_values("total_views", ascending=False)
    )
    creator_records = [
        {
            "Creator name": row["Creator name"],
            "videos": int(row["videos"]),
            "total_views": int(row["total_views"]),
            "paid_views": json_default(row["paid_views"]),
            "organic_views": int(row["organic_views"]),
            "likes": int(row["likes"]),
            "shares": int(row["shares"]),
            "engagement_rate": float(row["engagement_rate"]),
        }
        for _, row in creators.iterrows()
    ]

    top_videos = [
        row_to_dict(row)
        for _, row in videos_df.nlargest(5, "Total views").iterrows()
    ]

    history = [
        {
            "Date": row["Date"].isoformat(),
            "Total views": json_default(row["Total views"]),
            "Paid views": json_default(row["Paid views"]),
            "Organic views": json_default(row["Organic views"]),
        }
        for _, row in history_df.iterrows()
        if pd.notna(row["Date"])
    ]

    campaign_name = str(campaign.get("Campaign or link name", ""))
    resolved_period = period or infer_period(history_df)
    resolved_brand = brand or default_brand()
    resolved_source = source or infer_product_label(campaign_name, input_path)

    meta = {
        "n_creators": int(videos_df["Creator name"].nunique()),
        "n_videos": int(len(videos_df)),
        "period": resolved_period,
        "brand": resolved_brand,
        "source": resolved_source,
    }

    qa_notes = build_qa_notes(campaign, videos_df, history_df)

    return ReportContext(
        campaign=campaign,
        creators=creator_records,
        top_videos=top_videos,
        history=history,
        meta=meta,
        qa_notes=qa_notes,
        videos_df=videos_df,
        period=resolved_period,
        brand=resolved_brand,
        source=resolved_source,
        product_label=clean_product_label(resolved_source),
    )


def build_qa_notes(
    campaign: dict[str, Any],
    videos_df: pd.DataFrame,
    history_df: pd.DataFrame,
) -> list[str]:
    notes: list[str] = []
    campaign_views = int(campaign["Total views"])
    video_views = int(videos_df["Total views"].sum())
    gap = campaign_views - video_views
    if gap > 0:
        notes.append(
            f"Video-level view sum ({fmt_num(video_views)}) is lower than campaign total "
            f"({fmt_num(campaign_views)}) by {fmt_num(gap)} — use campaign totals for KPIs."
        )

    if not history_df.empty:
        paid_gt_total = history_df[
            history_df["Paid views"].fillna(0) > history_df["Total views"].fillna(0)
        ]
        if not paid_gt_total.empty:
            dates = pd.to_datetime(
                paid_gt_total["Date"].astype(str), format="%Y%m%d", errors="coerce"
            )
            for date in dates.dropna():
                notes.append(
                    f"Paid views exceeded total views on {date.strftime('%d %b %Y')} "
                    "(TikTok attribution lag — do not adjust)."
                )

    return notes


def write_data_json(ctx: ReportContext, output_path: Path) -> None:
    payload = {
        "campaign": ctx.campaign,
        "creators": ctx.creators,
        "top_videos": ctx.top_videos,
        "history": ctx.history,
        "meta": ctx.meta,
    }
    output_path.parent.mkdir(parents=True, exist_ok=True)
    with output_path.open("w", encoding="utf-8") as fh:
        json.dump(payload, fh, indent=2, ensure_ascii=False)
        fh.write("\n")


def delete_slide(prs: Presentation, index: int) -> None:
    slide_ids = prs.slides._sldIdLst
    slide_id = list(slide_ids)[index]
    prs.part.drop_rel(slide_id.rId)
    slide_ids.remove(slide_id)


def set_text_by_match(slide, match: str, new_text: str, *, font_size: int | None = None) -> bool:
    for shape in slide.shapes:
        if shape.has_text_frame and match in shape.text_frame.text:
            tf = shape.text_frame
            tf.clear()
            paragraph = tf.paragraphs[0]
            run = paragraph.add_run()
            run.text = new_text
            if font_size:
                run.font.size = Pt(font_size)
            return True
    return False


def add_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    *,
    size: int = 10,
    color: RGBColor | None = None,
    bold: bool = False,
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    paragraph = tf.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return box


def audience_summary(campaign: dict[str, Any]) -> tuple[str, str, str]:
    ages = parse_distribution(campaign.get("Age distribution"), "age")
    genders = parse_distribution(campaign.get("Gender distribution"), "gender")
    countries = parse_distribution(campaign.get("Countries or regions distribution"), "country")

    age_text = ", ".join(
        f"{item['age']} ({item['percentage'] * 100:.0f}%)" for item in ages[:2]
    )
    gender_text = " / ".join(
        f"{item['gender']} {item['percentage'] * 100:.0f}%"
        for item in genders
        if item.get("gender") in ("Male", "Female")
    )
    top_country = countries[0]["country"] if countries else "—"
    top_country_pct = countries[0]["percentage"] * 100 if countries else 0
    return age_text or "—", gender_text or "—", f"{top_country} {top_country_pct:.0f}%"


def spotlight_videos(videos_df: pd.DataFrame) -> tuple[pd.Series, pd.Series, pd.Series]:
    best_views = videos_df.nlargest(1, "Total views").iloc[0]
    best_shares = videos_df.nlargest(1, "Shares").iloc[0]
    filtered = videos_df[videos_df["Total views"] >= 5000]
    if filtered.empty:
        filtered = videos_df
    best_retention = filtered.nlargest(1, "Video completion rate").iloc[0]
    return best_views, best_shares, best_retention


def build_copy_blocks(ctx: ReportContext) -> dict[str, str]:
    camp = ctx.campaign
    videos = ctx.videos_df

    views = int(camp["Total views"])
    paid = int(camp["Paid views"])
    organic = int(camp["Organic views"])
    unique = int(camp["Unique viewers"])
    likes = int(camp["Likes"])
    comments = int(camp["Comments"])
    shares = int(camp["Shares"])
    favorites = int(camp["Favorites"])
    eng_rate = camp["Engagement rate"] * 100
    paid_eng = camp["Paid engagement rate"] * 100
    org_eng = camp["Organic engagement rate"] * 100
    avg_watch = camp["Average view time"]
    completion = camp["Video completion rate"] * 100
    sec2 = camp["2-Second video views"] * 100
    sec6 = camp["6-Second video views"] * 100
    cpm = camp["CPM"]
    cpc = camp["CPC"]
    cost = camp["Cost"]
    impressions = int(camp["Impressions"])

    age_text, gender_text, geo_text = audience_summary(camp)
    best_views, best_shares, best_retention = spotlight_videos(videos)
    best_eng = videos.nlargest(1, "Engagement rate").iloc[0]
    top_creator = ctx.creators[0] if ctx.creators else {"Creator name": "—", "total_views": 0}
    top_share = pct(top_creator["total_views"], views)
    low_videos = len(videos[videos["Total views"] < 500])

    overview_main = (
        f"Total views mencapai {fmt_num(views)}, dengan distribusi cukup seimbang:\n"
        f"Organic: {fmt_num(organic)} (≈{pct(organic, views)}%)\n"
        f"Paid: {fmt_num(paid)} (≈{pct(paid, views)}%)\n"
        f"Unique viewers: {fmt_num(unique)} → reach cukup luas dan tidak repetitif\n"
        f"Total impressions: {fmt_num(impressions)}\n"
        f"Insight: performance TTO Wound (Solution, Ointment, Bening) berjalan kuat — "
        f"paid amplification mendorong scale, organic tetap kontribusi {pct(organic, views)}%"
    )
    engagement_block = (
        f"Total likes: {fmt_num(likes)}\n"
        f"Comments: {comments:,}\n"
        f"Shares: {shares}\n"
        f"Favorites: {fmt_num(favorites)}\n"
        f"Overall engagement rate: {eng_rate:.2f}%\n"
        f"Paid engagement rate: {paid_eng:.2f}%\n"
        f"Organic engagement rate: {org_eng:.2f}%"
    )
    video_block = (
        f"Average watch time: {avg_watch:.1f} detik\n"
        f"Video completion rate: {completion:.1f}%\n"
        f"2-sec views: {sec2:.1f}%\n"
        f"6-sec views: {sec6:.1f}%\n"
        f"Cost: ${cost:.2f} | CPM: ${cpm:.2f} | CPC: ${cpc:.2f}"
    )
    meta_block = (
        f"Total Creators \t\t-> {ctx.meta['n_creators']} creators\n"
        f"Total Content videos\t-> {ctx.meta['n_videos']} videos\n"
        f"Period \t\t\t-> {ctx.period}\n"
        f"GSheets Recap \t\t-> link"
    )
    metrics = (
        f"Reach & Views: {fmt_num(views)} | Paid {fmt_num(paid)} / Organic {fmt_num(organic)}\n"
        f"Unique Viewers: {fmt_num(unique)} | Impressions: {fmt_num(impressions)}\n"
        f"Engagement: {eng_rate:.2f}% (Paid {paid_eng:.2f}% | Organic {org_eng:.2f}%)\n"
        f"Likes {fmt_num(likes)} | Comments {comments} | Shares {shares} | Favorites {fmt_num(favorites)}\n"
        f"Avg watch {avg_watch:.1f}s | Completion {completion:.1f}% | 2-sec {sec2:.1f}% → 6-sec {sec6:.1f}%\n"
        f"Cost ${cost:.2f} | CPM ${cpm:.2f} | CPC ${cpc:.2f} | Audience: {age_text}, {geo_text}"
    )
    metrics_detailed = (
        f"CAMPAIGN METRICS SUMMARY\n\n"
        f"Reach & Views\n• Total Views: {views:,} ({fmt_num(views)})\n"
        f"• Paid / Organic: {paid:,} / {organic:,}\n"
        f"• Unique Viewers: {unique:,}\n• Impressions: {impressions:,}\n\n"
        f"Engagement\n• Rate: {eng_rate:.2f}% (Paid {paid_eng:.2f}% | Organic {org_eng:.2f}%)\n"
        f"• Likes: {likes:,} | Comments: {comments} | Shares: {shares}\n\n"
        f"Video Performance\n• Avg Watch: {avg_watch:.1f}s | Completion: {completion:.2f}%\n"
        f"• 2-sec: {sec2:.1f}% → 6-sec: {sec6:.1f}% (drop {sec2 - sec6:.1f}pp)\n\n"
        f"Cost Efficiency\n• Total Cost: ${cost:.2f} | CPM: ${cpm:.2f} | CPC: ${cpc:.2f}\n\n"
        f"Audience\n• Age: {age_text}\n• Gender: {gender_text}\n• Geo: {geo_text}"
    )
    content_best = (
        f"Best Overall\n@{best_views['Creator name']}\n"
        f"{fmt_num(best_views['Total views'])} views\nLink to Video"
    )
    content_shared = (
        f"Most Shared\n@{best_shares['Creator name']}\n"
        f"{int(best_shares['Shares'])} shares\nLink to Video"
    )
    content_retention = (
        f"Best Retention Rate\n@{best_retention['Creator name']}\n"
        f"{best_retention['Video completion rate'] * 100:.1f}% completion\nLink to Video"
    )

    peak_day = ""
    if ctx.history:
        peak = max(ctx.history, key=lambda row: row.get("Total views") or 0)
        peak_date = datetime.fromisoformat(str(peak["Date"])).strftime("%d %b")
        peak_day = f"{peak_date}: ~{fmt_num(peak['Total views'])} views"

    worked = (
        "Strong visibility across Wound portfolio\n"
        "Exposure berhasil menjangkau audience secara luas via paid + organic\n\n"
        f"Cost efficiency is strong (CPM ${cpm:.2f})\n"
        "Campaign Wound berhasil scale dengan biaya relatif rendah\n\n"
        "Creator-led performance\n"
        f"@{top_creator['Creator name']} & @{best_eng['Creator name']} mendorong reach & engagement"
    )
    learnings = (
        "Paid media significantly lifts engagement\n"
        f"Paid ER {paid_eng:.2f}% vs organic {org_eng:.2f}% — ads meningkatkan interaksi\n\n"
        "Audience alignment on point\n"
        f"Core demo {age_text}, {geo_text} — relevan untuk Wound Care\n\n"
        f"Organic masih under-index ({pct(organic, views)}% views, {org_eng:.2f}% ER)\n"
        "Ada ruang growth konten organik across produk Wound"
    )
    improve = (
        "Engagement belum maksimal dibanding reach\n"
        "Content menjangkau banyak audience, interaksi (comments & shares) masih average\n\n"
        "Noticeable drop after first 2–6 seconds\n"
        f"{sec2:.0f}% pass 2-detik → hanya {sec6:.0f}% di 6-detik — storytelling perlu diperkuat\n\n"
        "Creator concentration risk\n"
        f"@{top_creator['Creator name']} kontribusi {top_share}% views — diversifikasi perlu diperkuat"
    )

    solution_worked = (
        "Strong paid reach at low cost\n"
        f"CPM ${cpm:.2f} — campaign scale dengan efisiensi biaya tinggi\n\n"
        "Creator efficiency terpusat\n"
        f"@{top_creator['Creator name']} mendorong {top_share}% total views via paid amplification\n\n"
        "Engagement spike on select content\n"
        f"@{best_eng['Creator name']} capai {best_eng['Engagement rate'] * 100:.1f}% engagement rate"
    )
    solution_learnings = (
        "Satu creator dominan performa\n"
        f"Konsentrasi views pada @{top_creator['Creator name']} — diversifikasi creator perlu dipertimbangkan\n\n"
        "Peak driven by paid burst\n"
        f"{peak_day or 'Paid burst'} — momentum bisa direplikasi\n\n"
        "Audience fit personal care\n"
        f"Core demo {age_text}, {geo_text} — selaras kategori antiseptik"
    )
    solution_improve = (
        f"Organic engagement lemah ({org_eng:.2f}%)\n"
        "Konten dilihat tapi belum memicu interaksi organik\n\n"
        "Retention drop setelah hook\n"
        f"{sec2:.0f}% pass 2-detik → hanya {sec6:.0f}% di 6-detik — storytelling perlu diperkuat\n\n"
        "Long-tail underperformance\n"
        f"{low_videos} video <500 views — evaluasi angle content"
    )

    return {
        "overview_main": overview_main,
        "engagement_block": engagement_block,
        "video_block": video_block,
        "meta_block": meta_block,
        "metrics": metrics,
        "metrics_detailed": metrics_detailed,
        "content_best": content_best,
        "content_shared": content_shared,
        "content_retention": content_retention,
        "worked": worked,
        "learnings": learnings,
        "improve": improve,
        "solution_worked": solution_worked,
        "solution_learnings": solution_learnings,
        "solution_improve": solution_improve,
    }


def populate_content_slide(slide, copy: dict[str, str]) -> None:
    set_text_by_match(slide, "Best Overall", copy["content_best"])
    set_text_by_match(slide, "Most Shared", copy["content_shared"])
    add_textbox(
        slide,
        6.8,
        1.5,
        2.8,
        2.5,
        copy["content_retention"],
        size=12,
        bold=True,
    )


def populate_overview_slide(slide, copy: dict[str, str]) -> None:
    set_text_by_match(slide, "Total views mencapai", copy["overview_main"])
    set_text_by_match(slide, "Total likes:", copy["engagement_block"])
    set_text_by_match(slide, "Average watch time:", copy["video_block"])


def populate_learnings_slide(slide, worked: str, learnings: str, improve: str) -> None:
    set_text_by_match(slide, "Strong visibility", worked)
    set_text_by_match(slide, "Audience alignment", learnings)
    set_text_by_match(slide, "Engagement belum", improve)
    add_textbox(slide, 3.5, 1.3, 2.8, 3.5, learnings, size=10)


def generate_wound_deck(
    ctx: ReportContext,
    template_path: Path,
    output_path: Path,
) -> None:
    if not template_path.exists():
        raise FileNotFoundError(f"Template not found: {template_path}")

    copy = build_copy_blocks(ctx)
    shutil.copy(template_path, output_path)
    prs = Presentation(str(output_path))

    set_text_by_match(
        prs.slides[0],
        "Betadine Wound",
        f"Betadine Wound\n({wound_title_subtitle(ctx.brand)})\nTTO Report & Learning",
    )

    populate_overview_slide(prs.slides[1], copy)

    for idx in (3, 7, 11):
        set_text_by_match(prs.slides[idx], "Total Creators", copy["meta_block"])

    for idx in (4, 8, 12):
        populate_content_slide(prs.slides[idx], copy)

    for idx, label in ((5, "SOLUTION"), (9, "OINTMENT"), (13, "BENING")):
        add_textbox(
            prs.slides[idx],
            0.4,
            0.9,
            9.2,
            4.2,
            f"{label} — Campaign Metrics\n{copy['metrics']}",
            size=10,
            color=GREEN,
        )

    populate_learnings_slide(
        prs.slides[14],
        copy["worked"],
        copy["learnings"],
        copy["improve"],
    )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


def generate_solution_deck(
    ctx: ReportContext,
    template_path: Path,
    output_path: Path,
) -> None:
    if not template_path.exists():
        raise FileNotFoundError(f"Template not found: {template_path}")

    copy = build_copy_blocks(ctx)
    shutil.copy(template_path, output_path)
    prs = Presentation(str(output_path))

    for idx in range(13, 5, -1):
        delete_slide(prs, idx)
    delete_slide(prs, 1)

    set_text_by_match(
        prs.slides[0],
        "Betadine Wound",
        f"{ctx.product_label}\nTTO Report & Learning",
    )

    slide_overview = prs.slides[2]
    set_text_by_match(slide_overview, "Total Creators", copy["meta_block"])
    if not any(
        "Total views mencapai" in shape.text_frame.text
        for shape in slide_overview.shapes
        if shape.has_text_frame
    ):
        add_textbox(
            slide_overview,
            0.4,
            1.0,
            9.2,
            4.0,
            "\n\n".join(
                [copy["overview_main"], copy["engagement_block"], copy["video_block"]]
            ),
            size=10,
        )
    else:
        set_text_by_match(slide_overview, "Total views mencapai", copy["overview_main"])
        set_text_by_match(slide_overview, "Total likes:", copy["engagement_block"])
        set_text_by_match(slide_overview, "Average watch time:", copy["video_block"])

    populate_content_slide(prs.slides[3], copy)
    add_textbox(
        prs.slides[4],
        0.4,
        0.9,
        9.2,
        4.2,
        copy["metrics_detailed"],
        size=10,
        color=GREEN,
    )
    populate_learnings_slide(
        prs.slides[5],
        copy["solution_worked"],
        copy["solution_learnings"],
        copy["solution_improve"],
    )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


def sanitize_filename(value: str) -> str:
    cleaned = re.sub(r'[<>:"/\\|?*]', "", value).strip()
    return re.sub(r"\s+", " ", cleaned)


def deck_filename(kind: str, ctx: ReportContext) -> str:
    period = sanitize_filename(ctx.period)
    if kind == "solution":
        return f"FOLO - Report TTO {sanitize_filename(ctx.product_label)} ({period}).pptx"
    return f"FOLO - Report TTO Betadine Wound ({period}).pptx"


def print_summary(ctx: ReportContext, outputs: list[Path]) -> None:
    print(f"Campaign: {ctx.campaign.get('Campaign or link name')}")
    print(f"Period: {ctx.period}")
    print(f"Brand: {ctx.brand}")
    print(f"Creators: {ctx.meta['n_creators']} | Videos: {ctx.meta['n_videos']}")
    print(f"Total views: {fmt_num(ctx.campaign['Total views'])}")
    if ctx.qa_notes:
        print("\nQA notes:")
        for note in ctx.qa_notes:
            print(f"  - {note}")
    print("\nOutputs:")
    for path in outputs:
        print(f"  - {path}")


def parse_args(argv: list[str] | None = None) -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Generate preview/data.json and FOLO-style PowerPoint decks from TikTok One Excel exports."
    )
    parser.add_argument("--input", required=True, type=Path, help="Path to TikTok One .xlsx export")
    parser.add_argument(
        "--format",
        choices=("solution", "wound", "both"),
        default="both",
        help="Deck format to generate (default: both)",
    )
    parser.add_argument("--period", help='Reporting period label, e.g. "Jun 2026"')
    parser.add_argument("--brand", help='Brand label for title/overview, e.g. "Solution, Ointment, Bening"')
    parser.add_argument("--source", help="Source label stored in data.json meta.source")
    parser.add_argument(
        "--template",
        type=Path,
        default=DEFAULT_TEMPLATE,
        help=f"Reference .pptx template (default: {DEFAULT_TEMPLATE})",
    )
    parser.add_argument(
        "--json-out",
        type=Path,
        default=DEFAULT_JSON,
        help=f"Output path for data.json (default: {DEFAULT_JSON})",
    )
    parser.add_argument(
        "--output-dir",
        type=Path,
        default=DEFAULT_OUTPUT_DIR,
        help=f"Directory for generated .pptx files (default: {DEFAULT_OUTPUT_DIR})",
    )
    parser.add_argument("--json-only", action="store_true", help="Only regenerate preview/data.json")
    parser.add_argument("--pptx-only", action="store_true", help="Only generate PowerPoint decks")
    return parser.parse_args(argv)


def main(argv: list[str] | None = None) -> int:
    args = parse_args(argv)
    if args.json_only and args.pptx_only:
        print("Use only one of --json-only or --pptx-only.", file=sys.stderr)
        return 2

    try:
        ctx = load_excel(
            args.input.resolve(),
            period=args.period,
            brand=args.brand,
            source=args.source,
        )
    except (FileNotFoundError, ValueError) as exc:
        print(str(exc), file=sys.stderr)
        return 1

    outputs: list[Path] = []

    if not args.pptx_only:
        write_data_json(ctx, args.json_out.resolve())
        outputs.append(args.json_out.resolve())

    if not args.json_only:
        output_dir = args.output_dir.resolve()
        template = args.template.resolve()
        if args.format in ("solution", "both"):
            solution_path = output_dir / deck_filename("solution", ctx)
            generate_solution_deck(ctx, template, solution_path)
            outputs.append(solution_path)
        if args.format in ("wound", "both"):
            wound_path = output_dir / deck_filename("wound", ctx)
            generate_wound_deck(ctx, template, wound_path)
            outputs.append(wound_path)

    print_summary(ctx, outputs)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
